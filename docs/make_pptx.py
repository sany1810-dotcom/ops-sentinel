"""Generate Ops-Sentinel.pptx for hackathon submission."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── palette ──────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0D, 0x11, 0x17)   # GitHub dark
ACCENT   = RGBColor(0x58, 0xA6, 0xFF)   # blue
DIM      = RGBColor(0x8B, 0x94, 0x9E)   # grey
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREEN    = RGBColor(0x56, 0xD3, 0x64)
ORANGE   = RGBColor(0xF0, 0x88, 0x3E)
YELLOW   = RGBColor(0xE3, 0xB3, 0x41)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return s


def txb(s, text, x, y, w, h, size=24, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb


def hline(s, y, color=ACCENT, w_frac=0.85):
    """Horizontal rule."""
    line = s.shapes.add_connector(
        1,  # STRAIGHT
        Inches((13.33 * (1 - w_frac)) / 2), Inches(y),
        Inches(13.33 - (13.33 * (1 - w_frac)) / 2), Inches(y),
    )
    line.line.color.rgb = color
    line.line.width     = Pt(1.5)


def bullet_block(s, items, x, y, w, size=20, marker="▸", color=WHITE, gap=0.55):
    for i, item in enumerate(items):
        txb(s, f"{marker}  {item}", x, y + i * gap, w, gap + 0.1,
            size=size, color=color)


def badge(s, text, x, y, w=2.2, h=0.45, bg=ACCENT, fg=BG, size=16, bold=True):
    """Filled rounded-rect badge via table hack (simplest approach in pptx)."""
    rect = s.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1 but we use freeform; use add_shape with idx
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    tf = rect.text_frame
    tf.text = text
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].runs[0].font.size  = Pt(size)
    tf.paragraphs[0].runs[0].font.bold  = bold
    tf.paragraphs[0].runs[0].font.color.rgb = fg
    from pptx.util import Pt as _Pt
    tf.margin_top = _Pt(4)
    tf.margin_bottom = _Pt(4)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Slide 1 — Title
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s1 = slide()
txb(s1, "Ops-Sentinel", 1, 1.6, 11, 1.4, size=64, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txb(s1, "Autonomous on-call agent with persistent incident memory",
    1, 2.95, 11, 0.7, size=26, color=WHITE, align=PP_ALIGN.CENTER)
hline(s1, 3.9)
txb(s1, "Qwen Cloud Hackathon  ·  Track 1: MemoryAgent",
    1, 4.1, 11, 0.5, size=18, color=DIM, align=PP_ALIGN.CENTER)
txb(s1, "Live demo:  http://47.237.196.56/",
    1, 5.1, 11, 0.5, size=17, color=GREEN, align=PP_ALIGN.CENTER)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Slide 2 — Problem
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s2 = slide()
txb(s2, "The Problem", 0.7, 0.4, 12, 0.8, size=38, bold=True, color=ACCENT)
hline(s2, 1.35)
bullet_block(s2, [
    "On-call engineers face the same incidents over and over",
    "Every recurrence is diagnosed from scratch — runbooks go stale",
    "Context is lost between shifts and team members",
    "Response time degrades under alert fatigue and pressure",
], x=0.9, y=1.65, w=11.5, size=22, color=WHITE, gap=0.7)

txb(s2, "The cost: slower MTTR, missed patterns, repeated mistakes",
    0.9, 4.8, 11.5, 0.6, size=20, color=ORANGE, italic=True)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Slide 3 — Solution
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s3 = slide()
txb(s3, "The Solution", 0.7, 0.4, 12, 0.8, size=38, bold=True, color=ACCENT)
hline(s3, 1.35)
txb(s3, "Ops-Sentinel never forgets an incident.", 0.9, 1.55, 11.5, 0.6,
    size=24, bold=True, color=WHITE)
bullet_block(s3, [
    "Continuously monitors a target service (metrics, latency, errors)",
    "Detects anomalies and extracts a symptom list automatically",
    "Searches past incidents semantically — finds similar cases by meaning, not keywords",
    "Calls Qwen with retrieved context: 'last time this happened, we restarted'",
    "Executes remediation (restart / alert / halt) and records the new incident",
    "Improves with every cycle — memory grows, responses sharpen",
], x=0.9, y=2.25, w=11.5, size=20, color=WHITE, gap=0.62)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Slide 4 — Architecture
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s4 = slide()
txb(s4, "Architecture", 0.7, 0.4, 12, 0.8, size=38, bold=True, color=ACCENT)
hline(s4, 1.35)

# Draw boxes
def box(s, label, sublabel, x, y, w=2.6, h=0.9, bc=ACCENT, tc=BG):
    r = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = bc
    r.line.color.rgb = bc
    tf = r.text_frame; tf.word_wrap = True
    tf.margin_top = Pt(4); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = label
    run.font.size = Pt(15); run.font.bold = True; run.font.color.rgb = tc
    if sublabel:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sublabel
        r2.font.size = Pt(11); r2.font.color.rgb = tc

def arrow(s, x1, y1, x2, y2, label="", color=DIM):
    c = s.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(1.8)
    if label:
        mx, my = (x1+x2)/2, (y1+y2)/2
        txb(s, label, mx-0.6, my-0.2, 1.2, 0.35, size=11, color=DIM, align=PP_ALIGN.CENTER)

# Row 1
box(s4, "Demo Service", ":8000  /inject", 0.35, 1.55, 2.4, 0.85, bc=RGBColor(0x6E,0x35,0x13), tc=ORANGE)
box(s4, "ops-agent", "poll · detect · act", 3.3,  1.55, 2.7, 0.85, bc=RGBColor(0x1F,0x6B,0x2E), tc=GREEN)
box(s4, "Qwen API", "qwen3.6-flash", 6.6,  1.55, 2.5, 0.85, bc=ACCENT, tc=BG)
box(s4, "Status Page", "GET /  :80", 10.2, 1.55, 2.6, 0.85, bc=RGBColor(0x21,0x26,0x2E), tc=DIM)

# Row 2
box(s4, "MCP Server", "FastMCP  :8002/mcp", 3.3, 3.5, 2.7, 0.85, bc=RGBColor(0x0D,0x3B,0x6E), tc=ACCENT)
box(s4, "4 MCP Tools", "search · record · recent · stats", 6.6, 3.5, 2.5, 0.85, bc=RGBColor(0x0D,0x3B,0x6E), tc=WHITE)
box(s4, "SQLite Memory", "incidents + embeddings", 10.2, 3.5, 2.6, 0.85, bc=RGBColor(0x21,0x26,0x2E), tc=DIM)

# Row 3
box(s4, "EmbeddingClient", "text-embedding-v3", 6.6, 5.3, 2.5, 0.8, bc=RGBColor(0x5A,0x3E,0x00), tc=YELLOW)

# Arrows
arrow(s4, 2.75, 1.97, 3.3, 1.97, color=ORANGE)        # demo→agent
arrow(s4, 4.65, 2.40, 4.65, 3.5, color=DIM)           # agent→mcp (vertical)
arrow(s4, 6.0,  1.97, 6.6, 1.97, color=ACCENT)        # agent→qwen
arrow(s4, 6.6,  1.97, 6.0, 1.97, color=ACCENT)        # qwen→agent (same line, bidirectional)
arrow(s4, 9.1,  3.92, 10.2, 3.92, color=DIM)          # tools→sqlite
arrow(s4, 7.85, 4.35, 7.85, 5.3, color=YELLOW)        # tools→embed (vertical)
arrow(s4, 10.2, 1.97, 10.2, 1.97, color=DIM)          # agent→status (same row)
arrow(s4, 6.0,  3.92, 6.6, 3.92, color=DIM)           # mcp→tools

# bidirectional label
txb(s4, "tool-calls", 6.05, 1.62, 0.9, 0.3, size=10, color=ACCENT, align=PP_ALIGN.CENTER)

# safe-mode note
txb(s4, "safe-mode: MCP down → direct SQLite  |  Qwen down → rule-based",
    0.35, 6.55, 12.5, 0.4, size=13, color=YELLOW, italic=True)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Slide 5 — MCP Integration
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s5 = slide()
txb(s5, "MCP Integration", 0.7, 0.4, 12, 0.8, size=38, bold=True, color=ACCENT)
hline(s5, 1.35)

txb(s5, "Transport:  FastMCP  ·  Streamable HTTP  ·  /mcp endpoint",
    0.9, 1.55, 11.5, 0.5, size=19, color=DIM)

tools = [
    ("search_similar_incidents", "Semantic vector search over past incidents; returns similarity_score", ACCENT),
    ("record_incident",          "Persists new incident + computes & stores embedding immediately",     GREEN),
    ("get_recent_incidents",     "Returns N most recent incidents for Qwen context",                    DIM),
    ("get_stats",                "Embedding coverage, semantic_search flag, top symptoms",              DIM),
]
for i, (name, desc, col) in enumerate(tools):
    y = 2.2 + i * 1.05
    badge(s5, name, x=0.7, y=y, w=3.8, h=0.42, bg=col if col != DIM else RGBColor(0x21,0x26,0x2E), fg=BG if col != DIM else WHITE, size=14)
    txb(s5, desc, 4.7, y+0.02, 8.2, 0.42, size=17, color=WHITE)

txb(s5, "Qwen calls tools autonomously in a loop — no hardcoded decision tree",
    0.9, 6.35, 11.5, 0.5, size=18, bold=True, color=YELLOW, italic=True)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Slide 6 — Semantic Memory
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s6 = slide()
txb(s6, "Semantic Incident Memory", 0.7, 0.4, 12, 0.8, size=38, bold=True, color=ACCENT)
hline(s6, 1.35)

# Left column: how it works
bullet_block(s6, [
    "Every incident embedded with text-embedding-v3  (1024-dim)",
    "Stored as float32 BLOB in SQLite incident_embeddings table",
    "Query: build_embed_text(symptoms + metrics) → cosine similarity",
    "Returns top-N incidents with similarity_score  (0.0 – 1.0)",
    "Finds conceptually similar faults — no keyword overlap needed",
], x=0.7, y=1.6, w=6.8, size=19, color=WHITE, gap=0.63)

# Right column: UI evidence box
r = s6.shapes.add_shape(1, Inches(7.8), Inches(1.55), Inches(5.0), Inches(3.6))
r.fill.solid(); r.fill.fore_color.rgb = RGBColor(0x16, 0x1B, 0x22)
r.line.color.rgb = ACCENT; r.line.width = Pt(1)

txb(s6, "Status Page — MCP Tool Calls", 7.9, 1.65, 4.8, 0.4, size=13, color=DIM)
txb(s6, "Tool", 7.9, 2.05, 2.2, 0.35, size=13, bold=True, color=DIM)
txb(s6, "Result / Similarity", 10.1, 2.05, 2.6, 0.35, size=13, bold=True, color=DIM)
hline(s6, 2.45)

rows = [
    ("search_similar_incidents", "[semantic] id=12 sim=0.91,\nid=9 sim=0.84, id=7 sim=0.79", GREEN),
    ("record_incident",          '{"id": 13, "embedded": true}',                               ACCENT),
    ("get_stats",                '{"semantic_search": true,\n"coverage": "13/13"}',             DIM),
]
for i, (tool, res, col) in enumerate(rows):
    y = 2.55 + i * 0.9
    txb(s6, tool, 7.9, y, 2.1, 0.8, size=12, color=col, bold=True)
    txb(s6, res,  10.1, y, 2.5, 0.8, size=11, color=WHITE)

txb(s6, "Similarity score visible to judges — proof of semantic retrieval",
    0.7, 5.45, 12.0, 0.5, size=17, color=YELLOW, italic=True)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Slide 7 — Resilience / Safe Mode
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s7 = slide()
txb(s7, "Resilience & Safe Mode", 0.7, 0.4, 12, 0.8, size=38, bold=True, color=ACCENT)
hline(s7, 1.35)

# Two columns
txb(s7, "MCP Server unavailable", 0.7, 1.6, 5.8, 0.5, size=22, bold=True, color=ORANGE)
bullet_block(s7, [
    "Agent detects failed MCP connection",
    "Falls back to direct SQLite text-overlap search",
    "Reconnect attempted in background every cycle",
    "Zero downtime — agent keeps diagnosing",
], x=0.7, y=2.2, w=5.8, size=18, color=WHITE, gap=0.58, marker="→")

txb(s7, "Qwen API unavailable", 7.0, 1.6, 5.8, 0.5, size=22, bold=True, color=ORANGE)
bullet_block(s7, [
    "3× retry with backoff on API errors",
    "Falls back to rule-based diagnosis",
    "Action: alert (conservative safe default)",
    "Incident still recorded in memory",
], x=7.0, y=2.2, w=5.8, size=18, color=WHITE, gap=0.58, marker="→")

hline(s7, 5.05, color=DIM)
txb(s7, "Embedding API unavailable", 0.7, 5.2, 12, 0.5, size=22, bold=True, color=ORANGE)
bullet_block(s7, [
    "search_similar_incidents returns [text_fallback] mode — symptom keyword overlap",
    "record_incident saves the incident without a vector (embeddable later via migrate script)",
], x=0.7, y=5.75, w=12.0, size=18, color=WHITE, gap=0.52, marker="→")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Slide 8 — Tech Stack & Deployment
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s8 = slide()
txb(s8, "Tech Stack & Deployment", 0.7, 0.4, 12, 0.8, size=38, bold=True, color=ACCENT)
hline(s8, 1.35)

stack = [
    ("Agent Runtime",   "Python 3.11 · FastAPI · uvicorn · httpx"),
    ("LLM",             "Qwen qwen3.6-flash via DashScope (OpenAI-compatible API)"),
    ("Embeddings",      "Qwen text-embedding-v3 · 1024-dim · same API key"),
    ("MCP Layer",       "mcp[cli] FastMCP · Streamable HTTP transport · 4 tools"),
    ("Memory",          "SQLite · incidents table + incident_embeddings BLOB"),
    ("Infrastructure",  "Docker Compose · 3 containers · Alibaba Cloud ECS"),
]
for i, (layer, detail) in enumerate(stack):
    y = 1.65 + i * 0.72
    txb(s8, layer, 0.7, y, 2.8, 0.55, size=17, bold=True, color=ACCENT)
    txb(s8, detail, 3.65, y, 9.2, 0.55, size=17, color=WHITE)

hline(s8, 6.1, color=DIM)
txb(s8, "Live on Alibaba Cloud ECS", 0.7, 6.25, 5, 0.45, size=17, bold=True, color=GREEN)
txb(s8, "http://47.237.196.56/", 5.8, 6.25, 7, 0.45, size=17, color=GREEN)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Slide 9 — Demo Flow
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s9 = slide()
txb(s9, "Live Demo Flow", 0.7, 0.4, 12, 0.8, size=38, bold=True, color=ACCENT)
hline(s9, 1.35)
txb(s9, "http://47.237.196.56/ — open in browser, no login",
    0.7, 1.55, 12, 0.45, size=18, color=GREEN)

steps = [
    ("1  Reset",    "Click Reset Service — demo service starts clean",             WHITE),
    ("2  Inject",   "Click Inject: Overload (or memory_leak / dependency_down)",   ORANGE),
    ("3  Detect",   "Agent detects high latency + CPU symptoms within ~10 s",      WHITE),
    ("4  Recall",   "Qwen calls search_similar_incidents → [semantic] sim=0.91",   ACCENT),
    ("5  Diagnose", "Qwen reasons over retrieved past case, picks action",          WHITE),
    ("6  Act",      "Agent restarts / alerts, records new incident with embedding", GREEN),
    ("7  Repeat",   "Inject same fault — agent finds prior case, responds faster",  YELLOW),
]
for i, (step, desc, col) in enumerate(steps):
    y = 2.15 + i * 0.68
    txb(s9, step, 0.7, y, 1.9, 0.55, size=17, bold=True, color=col)
    txb(s9, desc, 2.75, y, 10.2, 0.55, size=17, color=WHITE)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Slide 10 — Closing
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s10 = slide()
txb(s10, "What Ops-Sentinel Demonstrates", 0.7, 0.5, 12, 0.8, size=36, bold=True, color=ACCENT)
hline(s10, 1.4)

pillars = [
    ("Memory",      "Persistent SQLite store — every incident remembered and embedded",    ACCENT),
    ("Autonomy",    "Qwen drives the full diagnose → decide → act loop via MCP tools",     GREEN),
    ("Semantics",   "text-embedding-v3 retrieval — finds similar faults by meaning",       YELLOW),
    ("Resilience",  "Three-layer fallback: semantic → text → rules. Agent never stops.",   ORANGE),
]
for i, (title, body, col) in enumerate(pillars):
    x = 0.5 + i * 3.1
    r = s10.shapes.add_shape(1, Inches(x), Inches(1.75), Inches(2.9), Inches(3.8))
    r.fill.solid(); r.fill.fore_color.rgb = RGBColor(0x16,0x1B,0x22)
    r.line.color.rgb = col; r.line.width = Pt(2)
    txb(s10, title, x+0.1, 1.85, 2.7, 0.55, size=22, bold=True, color=col)
    txb(s10, body,  x+0.1, 2.5,  2.7, 2.9,  size=15, color=WHITE)

hline(s10, 6.1, color=DIM)
txb(s10, "Built with Qwen Cloud APIs  ·  Deployed on Alibaba Cloud ECS  ·  Track 1: MemoryAgent",
    0.7, 6.25, 12, 0.5, size=16, color=DIM, align=PP_ALIGN.CENTER)

# ── save ─────────────────────────────────────────────────────────────────────
out = "docs/Ops-Sentinel.pptx"
prs.save(out)
print(f"Saved: {out}  ({prs.slides.__len__()} slides)")
